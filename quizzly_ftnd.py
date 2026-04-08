import os
import streamlit as st
import tempfile
import time
import traceback
from openai import OpenAIError
import mimetypes
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from PIL import Image
import uuid
import requests
from bs4 import BeautifulSoup

# Cleaned imports: no file conversion dependencies needed
from quizzly_bknd_gnrt import setup_api, get_page_count, create_extraction_chain, create_generation_chain
from quizzly_bknd_vrf import verify_quiz

# ~2500 characters of plain text ≈ one page for quiz sizing (web sources have no PDF page count)
WEB_CHARS_PER_PAGE = 2500
WEB_TEXT_PER_URL_CAP = 12000
# Each successful URL counts at least this many pseudo-pages for max-question math so one article can reach min 3 questions (ceil(6/2)).
WEB_PAGES_FLOOR_PER_URL = 6
MIN_QUESTIONS = 3


def extract_readable_text(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    parts = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li"]):
        t = el.get_text(" ", strip=True)
        if t:
            parts.append(t)

    text = "\n".join(parts).strip()

    if len(text) < 200:
        text = soup.get_text(separator="\n", strip=True)
        text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])

    return text


def fetch_website_text(url: str) -> tuple[bool, str]:
    """Fetch one URL and return (ok, extracted_text). Text is capped for downstream use."""
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    }

    resp = requests.get(url, headers=headers, timeout=20, allow_redirects=True)
    resp.raise_for_status()

    text = extract_readable_text(resp.text)

    if len(text) < 250:
        if url.startswith("https://"):
            fallback_url = "https://r.jina.ai/https://" + url[len("https://"):]
        elif url.startswith("http://"):
            fallback_url = "https://r.jina.ai/http://" + url[len("http://"):]
        else:
            fallback_url = "https://r.jina.ai/http://" + url

        fb = requests.get(fallback_url, headers=headers, timeout=20, allow_redirects=True)
        fb.raise_for_status()
        text = extract_readable_text(fb.text)

    if len(text) < 250:
        return False, ""

    return True, text[:WEB_TEXT_PER_URL_CAP]


def pseudo_pages_from_web_text(text: str) -> int:
    """Map extracted web text to a page count for max-question math (same formula as files: max_q ≈ pages//2)."""
    if not text:
        return 0
    return max(1, len(text) // WEB_CHARS_PER_PAGE)


def docx_to_pdf(input_path):
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")
    
    doc = Document(input_path)
    c = canvas.Canvas(output_path)

    y = 800
    for para in doc.paragraphs:
        text = para.text
        if y < 50:
            c.showPage()
            y = 800
        c.drawString(50, y, text[:100])
        y -= 15

    c.save()
    return output_path

def pptx_to_pdf(input_path):
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")
    
    prs = Presentation(input_path)
    c = canvas.Canvas(output_path)

    for slide in prs.slides:
        y = 800
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if y < 50:
                    c.showPage()
                    y = 800
                c.drawString(50, y, text[:100])
                y -= 15
        c.showPage()

    c.save()
    return output_path

def image_to_pdf(input_path):
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")
    
    image = Image.open(input_path).convert("RGB")
    image.save(output_path, "PDF")
    
    return output_path


st.set_page_config(page_title="Quizzly", page_icon="📖", layout="wide")

# Initialize Session States for stateful UI
if 'quiz_data' not in st.session_state:
    st.session_state.quiz_data = None
if 'error_notebook' not in st.session_state:
    st.session_state.error_notebook = []
if 'verification_report' not in st.session_state:
    st.session_state.verification_report = None
if 'generation_time' not in st.session_state:
    st.session_state.generation_time = None
if 'current_paths' not in st.session_state:
    st.session_state.current_paths = []
if "workflow_status_label" not in st.session_state:
    st.session_state.workflow_status_label = None
if "workflow_status_lines" not in st.session_state:
    st.session_state.workflow_status_lines = []

def main():
    st.title("📖 Quizzly: Automated Quiz Generator")
    st.markdown("Transform passive reading into active mastery. Upload documents or links to generate a verified, targeted quiz based on Bloom's Taxonomy.")

    # API Key Check via Streamlit Secrets
    if "OPENAI_API_KEY" not in st.secrets:
        st.error("⚠️ Please set the OPENAI_API_KEY in the Streamlit secrets.")
        return
    
    # Set it as an environment variable so LangChain and OpenAI clients pick it up automatically
    os.environ["OPENAI_API_KEY"] = st.secrets["OPENAI_API_KEY"]

    # --- Sidebar: Upload & Settings ---
    with st.sidebar:
        st.header("Study Materials")

        source_mode = st.radio(
            "Material source (choose one)",
            ["Upload files", "Website links"],
            horizontal=True,
            help="Use either uploaded files (up to 5) or website URLs (up to 3), not both.",
        )

        uploaded_files = None
        website_urls: list[str] = []

        if source_mode == "Upload files":
            uploaded_files = st.file_uploader(
                "Upload files (PDF, DOCX, PPTX, TXT, PNG) — max 5",
                type=["pdf", "docx", "pptx", "txt", "png"],
                accept_multiple_files=True,
            )
            if uploaded_files and len(uploaded_files) > 5:
                st.error("Please upload at most 5 files. Remove extras and try again.")
                uploaded_files = None
        else:
            st.caption("Enter up to 3 URLs (one per field). Google search pages often fail; use article links.")
            u1 = st.text_input("Website URL 1", key="web_url_1")
            u2 = st.text_input("Website URL 2", key="web_url_2")
            u3 = st.text_input("Website URL 3", key="web_url_3")
            website_urls = [u.strip() for u in (u1, u2, u3) if u and u.strip()]
            if len(website_urls) > 3:
                website_urls = website_urls[:3]

        num_questions = MIN_QUESTIONS
        generate_btn = False

        has_files = bool(uploaded_files)
        has_urls = bool(website_urls)

        if (source_mode == "Upload files" and has_files) or (source_mode == "Website links" and has_urls):
            temp_dir = tempfile.gettempdir()
            processed_paths: list[str] = []
            total_pages = 0
            web_blocks: list[tuple[str, str]] = []
            web_fetch_errors: list[str] = []
            source_count = 0
            files_eligible = False
            web_eligible = False

            if source_mode == "Upload files":
                st.session_state.web_text = ""
                for uf in uploaded_files:
                    temp_path = os.path.join(temp_dir, uf.name)
                    with open(temp_path, "wb") as f:
                        f.write(uf.getbuffer())

                    ext = os.path.splitext(temp_path)[1].lower()

                    if ext == ".docx":
                        new_path = docx_to_pdf(temp_path)
                    elif ext == ".pptx":
                        new_path = pptx_to_pdf(temp_path)
                    elif ext in [".png", ".jpg", ".jpeg"]:
                        new_path = image_to_pdf(temp_path)
                    else:
                        new_path = temp_path

                    processed_paths.append(new_path)
                    total_pages += get_page_count(new_path)

                source_count = len(processed_paths)
                raw_max = total_pages // 2
                files_eligible = raw_max >= MIN_QUESTIONS
                max_questions = raw_max if files_eligible else MIN_QUESTIONS
                st.caption(
                    f"Max questions = total pages ÷ 2 (minimum quiz size is {MIN_QUESTIONS}). "
                    f"Your files: **{total_pages}** page(s) → cap **{raw_max}**."
                )
                if not files_eligible:
                    st.warning(
                        f"Need at least {MIN_QUESTIONS * 2} total pages to allow a {MIN_QUESTIONS}-question quiz "
                        f"(current cap from pages ÷ 2 is {raw_max}). Add more files or pages."
                    )

            else:
                st.session_state.current_paths = []
                for url in website_urls:
                    try:
                        ok, chunk = fetch_website_text(url)
                        if ok and chunk:
                            web_blocks.append((url, chunk))
                        else:
                            web_fetch_errors.append(
                                f"{url}: too little readable text (try a direct article, not a search results page)."
                            )
                    except Exception as e:
                        web_fetch_errors.append(f"{url}: {e}")

                if web_fetch_errors:
                    for msg in web_fetch_errors:
                        st.error(msg)

                source_count = len(web_blocks)
                total_web_pages = sum(
                    max(WEB_PAGES_FLOOR_PER_URL, pseudo_pages_from_web_text(t)) for _, t in web_blocks
                )
                if not web_blocks:
                    total_web_pages = 0

                combined_web = "\n\n---\n\n".join(
                    f"Source: {u}\n\n{t}" for u, t in web_blocks
                )
                st.session_state.web_text = combined_web

                raw_max = total_web_pages // 2
                web_eligible = raw_max >= MIN_QUESTIONS
                max_questions = raw_max if web_eligible else MIN_QUESTIONS
                st.caption(
                    f"Max questions = estimated pages ÷ 2 (same idea as files). "
                    f"Each link counts at least **{WEB_PAGES_FLOOR_PER_URL}** pseudo-pages, plus more for long text "
                    f"(~{WEB_CHARS_PER_PAGE} characters ≈ 1 page). "
                    f"Estimated **{total_web_pages}** page(s) from **{source_count}** link(s) → cap **{raw_max}**."
                )
                if web_blocks and not web_eligible:
                    st.warning(
                        f"Fetched text is too short for a {MIN_QUESTIONS}-question cap (effective pages ÷ 2 = {raw_max}). "
                        "Try longer articles or more URLs."
                    )

            st.session_state.current_paths = processed_paths

            if source_count == 0:
                st.warning("Materials loaded: 0 sources — check your URLs or switch to file upload.")
            else:
                st.success(f"Materials Loaded: {source_count} source(s) detected.")
            if source_mode == "Upload files":
                eligible = files_eligible
            else:
                eligible = bool(web_blocks) and web_eligible

            if eligible:
                st.info(f"To maintain context quality, max questions is set to {max_questions}.")

            st.header("Quiz Settings")
            if eligible:
                num_questions = st.number_input(
                    "Number of Questions",
                    min_value=MIN_QUESTIONS,
                    max_value=max_questions,
                    value=MIN_QUESTIONS,
                )
            else:
                st.caption("Not enough material for the minimum quiz size; increase pages or web text to choose question count.")
                num_questions = MIN_QUESTIONS

            st.session_state.current_paths = processed_paths
            if source_mode == "Upload files":
                can_generate = bool(processed_paths) and files_eligible
            else:
                can_generate = bool(web_blocks) and web_eligible

            if not can_generate:
                st.info("Add valid material (files or fetchable URLs) to generate a quiz.")
            generate_btn = st.button("Generate & Verify Quiz", type="primary", disabled=(not can_generate))

    # --- Main Area: Processing & Display ---
    

    # --- View 1: Quiz Execution ---
    col1, col2 = st.columns([2, 1], gap="large")
    
    with col1:
        if st.session_state.workflow_status_label:
            st.status(st.session_state.workflow_status_label, state="complete", expanded=False)
            if st.session_state.workflow_status_lines:
                with st.expander("View workflow details"):
                    st.code("\n".join(st.session_state.workflow_status_lines))

        if generate_btn:
            st.session_state.workflow_status_label = None
            st.session_state.workflow_status_lines = []

            def log_line(s: str):
                st.session_state.workflow_status_lines.append(s)
                st.write(s)
            # 1. Start the timer right as the button is clicked
            start_time = time.time() 

            with st.status("Processing Document Workflow...", expanded=True) as status:
                try:
                    client = setup_api()
                    
                    log_line("Uploading to secure environment...")
                    oai_file_ids = []
                    mime_types = {
                        ".pdf": "application/pdf",
                        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        ".txt": "text/plain",
                        ".png": "image/png"
                    }

                    for fp in st.session_state.current_paths:
                        mime_type, _ = mimetypes.guess_type(fp)
                        if mime_type is None:
                            mime_type = "application/octet-stream"

                        st.write(f"Uploading: {fp} | MIME: {mime_type}")  # debug line

                        with open(fp, "rb") as f:
                            oai_file = client.files.create(
                                file=(os.path.basename(fp), f, mime_type),
                                purpose="user_data"
                            )

                        oai_file_ids.append(oai_file.id)
                    
                    log_line("Extracting core concepts...")
                    extractor = create_extraction_chain()
                    # Pass both file_ids and the direct web_text
                    concepts_resp = extractor.invoke({
                        "file_ids": oai_file_ids, 
                        "web_context": st.session_state.get("web_text", "")
                    })
                    concepts = concepts_resp.get("concepts") or []
                    if not concepts:
                        raise ValueError("Failed to extract concepts from the provided materials (website may be unreadable).")

                    log_line(f"Generating {num_questions} questions...")
                    generator = create_generation_chain(num_questions)
                    quiz_data = generator.invoke({
                        "file_ids": oai_file_ids,
                        "concepts_list": ", ".join(concepts),
                        "web_context": st.session_state.get("web_text", "")
                    })
                    
                    log_line("Running quiz verification checks...")
                    report = verify_quiz(concepts, quiz_data, num_questions)
                    
                    st.session_state.verification_report = report
                    st.session_state.quiz_data = quiz_data
                    
                    # Cleanup
                    for fp in st.session_state.current_paths:
                        try:
                            os.remove(fp)
                        except Exception:
                            pass
                    
                    # 3. Calculate elapsed time and update the status label dynamically
                    elapsed_time = time.time() - start_time
                    st.session_state.generation_time = elapsed_time
                    st.session_state.workflow_status_label = f"Workflow Complete in {elapsed_time:.1f} secs!"
                    status.update(label=st.session_state.workflow_status_label, state="complete", expanded=False)
                    
                except OpenAIError as e:
                    # Catches issues specifically related to OpenAI (Auth, Rate Limits, Timeouts)
                    status.update(label="OpenAI API Error", state="error")
                    st.error("There was a problem communicating with OpenAI. Check your API key, billing limits, or network connection.")
                    st.info(f"**Details:** {str(e)}")
                    
                except ValueError as e:
                    # Catches missing environment variables or bad inputs
                    status.update(label="Configuration Error", state="error")
                    st.error("A configuration or input value error occurred.")
                    st.info(f"**Details:** {str(e)}")
                    
                except Exception as e:
                    # The fallback for any other unexpected Python or LangChain errors
                    error_type = type(e).__name__
                    status.update(label=f"System Error: {error_type}", state="error")
                    st.error(f"The workflow failed due to an unexpected {error_type}.")
                    st.info(f"**Details:** {str(e)}")
                    
                    # Hidden expander for developers to see the exact line number of the crash
                    with st.expander("🛠️ Show Detailed Stack Trace (For Debugging)"):
                        st.code(traceback.format_exc(), language="python")

        if st.session_state.quiz_data:
            # Show verification results in an expander
            if st.session_state.verification_report:
                report = st.session_state.verification_report
                with st.expander("🛠️ View Comprehensive Verification Report"):
                    
                    # 1. Unpack all 6 metrics
                    passed = report.get('passed_constraints', 'Unknown')
                    c_score = report.get('constraint_score', 0.0)
                    c_feedback = report.get('constraint_feedback', [])
                    f_score = report.get('fidelity_score', 'N/A')
                    p_score = report.get('pedagogical_score', 'N/A')
                    reasoning = report.get('evaluator_reasoning', 'No reasoning provided by evaluator.')
                    
                    # 2. Render Overall Status
                    status_icon = "✅ PASSED" if passed else "❌ FAILED"
                    st.subheader(f"Pipeline Status: {status_icon}")
                    st.divider()
                    
                    # 3. Render Code-Based Metrics (Unit Tests)
                    st.markdown("#### 1. Structural Constraints (Code-Based Grading)")
                    st.write(f"**Constraint Score:** {c_score * 100}%")
                    for feedback_item in c_feedback:
                        # Use green checkmarks for passes, red X's for fails
                        icon = "✅" if "Pass" in feedback_item else "❌"
                        st.markdown(f"{icon} {feedback_item}")
                        
                    st.divider()
                    
                    # 4. Render LLM-Based Metrics (Integration Tests)
                    st.markdown("#### 2. Quality Evaluation (LLM-Based Grading)")
                    colA, colB = st.columns(2)
                    with colA:
                        st.metric(label="Task Fidelity Score", value=f"{f_score}/5")
                    with colB:
                        st.metric(label="Pedagogical Score", value=f"{p_score}/5")
                    
                    st.markdown("**Evaluator Reasoning:**")
                    st.info(reasoning)

            st.divider()
            st.subheader(st.session_state.quiz_data.get("quiz_title", "Assessment"))
            
            # Interactive Quiz
            with st.form("quiz_form"):
                user_answers = {}
                for i, q in enumerate(st.session_state.quiz_data.get("questions", [])):
                    # Add the Difficulty badge to the question header
                    difficulty = q.get('difficulty', 'Unrated')
                    st.markdown(f"**{i+1}. {q['question_text']}** *(Difficulty: {difficulty})*")
                    
                    user_answers[q['id']] = st.radio(
                        "Select an option:", 
                        q['options'], 
                        key=f"q_{q['id']}", 
                        index=None
                    )
                    st.write("---")
                
                submitted = st.form_submit_button("Submit Answers")
                
                if submitted:
                    for q in st.session_state.quiz_data["questions"]:
                        user_ans = user_answers[q['id']]
                        
                        if not user_ans:
                            st.warning(f"Question {q['id']} was left blank.")
                            continue
                            
                        user_letter = user_ans[0] # Extracts "A", "B", etc.
                        
                        # Format the explanation so Markdown renders the newlines perfectly
                        formatted_explanation = q['explanation'].replace('\n', '\n\n')
                        
                        if user_letter == q['correct_option']:
                            st.success(f"**Q{q['id']}:** Correct! ✅")
                            # Optionally show the explanation even when correct, for reinforcement!
                            with st.expander("Show detailed explanation"):
                                st.markdown(formatted_explanation)
                        else:
                            st.error(f"**Q{q['id']}:** Incorrect. The answer is {q['correct_option']}.")
                            st.info(formatted_explanation)
                            
                            # Add to Error Notebook if not already there
                            error_entry = {
                                "question": q['question_text'],
                                "user_wrong": user_ans,
                                "explanation": formatted_explanation
                            }
                            if error_entry not in st.session_state.error_notebook:
                                st.session_state.error_notebook.append(error_entry)

    # --- View 2: Error Notebook Right Panel ---
    with col2:
        # Wrap the entire right section in a bordered container with a set height
        with st.container(border=True, height=750):
            st.header("📓 Error Notebook")
            st.markdown("Review your mistakes to reinforce learning.")
            st.divider()

            if not st.session_state.error_notebook:
                st.info("No errors logged yet. Great job!")
            else:
                for idx, error in enumerate(st.session_state.error_notebook):
                    with st.expander(f"Review Question {idx + 1}"):
                        st.markdown(f"**Q:** {error['question']}")
                        st.markdown(f"❌ *You answered: {error['user_wrong']}*")
                        # Using markdown here guarantees the \n\n breaks render nicely
                        st.markdown(f"💡 **Correction:**\n\n{error['explanation']}")
                
                st.divider()
                # Adding use_container_width=True makes the button span the panel nicely
                if st.button("Clear Notebook", use_container_width=True):
                    st.session_state.error_notebook = []
                    st.rerun()

if __name__ == "__main__":
    main()
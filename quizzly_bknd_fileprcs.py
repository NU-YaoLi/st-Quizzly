"""File conversion and web text extraction for Quizzly (no Streamlit)."""

import ipaddress
import os
import socket
import tempfile
import uuid
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image
from pptx import Presentation
from reportlab.pdfgen import canvas

from quizzly_config import WEB_CHARS_PER_PAGE, WEB_TEXT_PER_URL_CAP


def _is_safe_http_url(url: str) -> bool:
    """Basic SSRF guard: allow only http/https and block localhost/private IP literals."""
    try:
        p = urlparse(url)
    except Exception:
        return False

    if p.scheme not in {"http", "https"}:
        return False
    if not p.netloc:
        return False

    hostname = (p.hostname or "").strip().lower()
    if not hostname:
        return False

    if hostname in {"localhost"} or hostname.endswith(".localhost") or hostname.endswith(".local"):
        return False

    # If hostname is an IP literal, block private/internal ranges.
    try:
        ip = ipaddress.ip_address(hostname)
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            return False
    except ValueError:
        # Not an IP literal; allow (DNS-level blocking is out of scope here).
        pass

    # DNS-based block: if hostname resolves to internal IPs, block.
    try:
        infos = socket.getaddrinfo(hostname, p.port or (443 if p.scheme == "https" else 80), type=socket.SOCK_STREAM)
        for info in infos:
            addr = info[4][0]
            try:
                ip = ipaddress.ip_address(addr)
                if (
                    ip.is_private
                    or ip.is_loopback
                    or ip.is_link_local
                    or ip.is_multicast
                    or ip.is_reserved
                    or ip.is_unspecified
                ):
                    return False
            except ValueError:
                continue
    except Exception:
        # If DNS fails, treat as unsafe.
        return False

    return True


def extract_readable_text(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    parts = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li"]):
        t = el.get_text(" ", strip=True)
        if t:
            parts.append(t)

    text = "\n".join(parts).strip()

    if len(text) < 200:
        text = soup.get_text(separator="\n", strip=True)
        text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])

    return text


def fetch_website_text(url: str) -> tuple[bool, str]:
    """Fetch one URL and return (ok, extracted_text). Text is capped for downstream use."""
    if not _is_safe_http_url(url):
        return False, ""

    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    }

    resp = requests.get(url, headers=headers, timeout=(10, 20), allow_redirects=True)
    resp.raise_for_status()

    text = extract_readable_text(resp.text)

    if len(text) < 250:
        if url.startswith("https://"):
            fallback_url = "https://r.jina.ai/https://" + url[len("https://"):]
        elif url.startswith("http://"):
            fallback_url = "https://r.jina.ai/http://" + url[len("http://"):]
        else:
            fallback_url = "https://r.jina.ai/http://" + url

        if not _is_safe_http_url(fallback_url):
            return False, ""

        fb = requests.get(fallback_url, headers=headers, timeout=(10, 20), allow_redirects=True)
        fb.raise_for_status()
        text = extract_readable_text(fb.text)

    if len(text) < 250:
        return False, ""

    return True, text[:WEB_TEXT_PER_URL_CAP]


def pseudo_pages_from_web_text(text: str) -> int:
    """Map extracted web text to a page count for max-question sizing."""
    if not text:
        return 0
    return max(1, len(text) // WEB_CHARS_PER_PAGE)


def docx_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    doc = Document(input_path)
    c = canvas.Canvas(output_path)

    y = 800
    for para in doc.paragraphs:
        text = para.text
        if y < 50:
            c.showPage()
            y = 800
        c.drawString(50, y, text[:100])
        y -= 15

    c.save()
    return output_path


def pptx_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    prs = Presentation(input_path)
    c = canvas.Canvas(output_path)

    for slide in prs.slides:
        y = 800
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if y < 50:
                    c.showPage()
                    y = 800
                c.drawString(50, y, text[:100])
                y -= 15
        c.showPage()

    c.save()
    return output_path


def image_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    image = Image.open(input_path).convert("RGB")
    image.save(output_path, "PDF")

    return output_path

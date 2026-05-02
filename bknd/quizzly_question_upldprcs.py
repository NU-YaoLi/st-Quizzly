"""File conversion and web text extraction for Quizzly (no Streamlit UI)."""

import ipaddress
import os
import socket
import tempfile
import uuid
from urllib.parse import urlparse

import requests
import streamlit as st
from bs4 import BeautifulSoup
from docx import Document
from PIL import Image
from pptx import Presentation
from reportlab.pdfgen import canvas

from quizzly_config import MAX_WEB_URL_SLOTS, WEB_CHARS_PER_PAGE, WEB_TEXT_PER_URL_CAP

PENDING_REMOVE_URL_INDEX = "_pending_remove_url_index"
MAX_REDIRECTS = 5
MAX_WEB_RESPONSE_BYTES = 2_000_000  # 2MB cap to prevent huge-page DoS

# Defensive: prevent decompression-bomb style images from consuming huge RAM.
# (Pillow raises DecompressionBombError / warnings when exceeded.)
Image.MAX_IMAGE_PIXELS = 20_000_000


def apply_pending_web_url_removal() -> None:
    """
    Apply a URL row removal before URL widgets mount.

    This avoids Streamlit widget state errors by shifting values in session_state
    before rendering the text_input widgets.
    """
    pending = st.session_state.pop(PENDING_REMOVE_URL_INDEX, None)
    if pending is None:
        return
    n = min(int(st.session_state.get("web_url_slot_count", 1)), MAX_WEB_URL_SLOTS)
    if n <= 1 or pending < 0 or pending >= n:
        return

    vals = [str(st.session_state.get(f"web_url_{i}", "") or "") for i in range(n)]
    new_vals = vals[:pending] + vals[pending + 1 :]

    for k in list(st.session_state.keys()):
        if isinstance(k, str) and k.startswith("web_url_") and k[8:].isdigit():
            del st.session_state[k]

    st.session_state["web_url_slot_count"] = max(1, len(new_vals))
    for i, v in enumerate(new_vals):
        st.session_state[f"web_url_{i}"] = v


def _check_http_url_safety(url: str) -> tuple[bool, str]:
    """SSRF guard: allow only http/https and block localhost/private IPs (including DNS-resolved)."""
    try:
        p = urlparse(url)
    except Exception:
        return False, "invalid_url"

    if p.scheme not in {"http", "https"}:
        return False, "invalid_scheme"
    if not p.netloc:
        return False, "missing_host"

    hostname = (p.hostname or "").strip().lower()
    if not hostname:
        return False, "missing_host"

    if hostname in {"localhost"} or hostname.endswith(".localhost") or hostname.endswith(".local"):
        return False, "blocked_localhost"

    # If hostname is an IP literal, block private/internal ranges.
    try:
        ip = ipaddress.ip_address(hostname)
        if (
            ip.is_private
            or ip.is_loopback
            or ip.is_link_local
            or ip.is_multicast
            or ip.is_reserved
            or ip.is_unspecified
        ):
            return False, "blocked_private_ip"
    except ValueError:
        # Not an IP literal; continue with DNS checks.
        pass

    # DNS-based block: if hostname resolves to internal IPs, block.
    try:
        infos = socket.getaddrinfo(
            hostname,
            p.port or (443 if p.scheme == "https" else 80),
            type=socket.SOCK_STREAM,
        )
        for info in infos:
            addr = info[4][0]
            try:
                ip = ipaddress.ip_address(addr)
                if (
                    ip.is_private
                    or ip.is_loopback
                    or ip.is_link_local
                    or ip.is_multicast
                    or ip.is_reserved
                    or ip.is_unspecified
                ):
                    return False, "blocked_private_ip"
            except ValueError:
                continue
    except Exception:
        # If DNS fails, treat as unsafe.
        return False, "dns_failed"

    return True, "ok"


def extract_readable_text(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")

    for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        tag.decompose()

    parts = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li"]):
        t = el.get_text(" ", strip=True)
        if t:
            parts.append(t)

    text = "\n".join(parts).strip()

    if len(text) < 200:
        text = soup.get_text(separator="\n", strip=True)
        text = "\n".join([line.strip() for line in text.splitlines() if line.strip()])

    return text


def fetch_website_text(url: str) -> tuple[bool, str, str]:
    """Fetch one URL and return (ok, extracted_text, reason). Text is capped for downstream use."""
    headers = {
        "User-Agent": "Mozilla/5.0",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    }

    # SSRF-safe fetch: validate each redirect hop (and disallow proxy fallbacks).
    session = requests.Session()
    # Avoid surprising behavior from environment proxies (security + reproducibility).
    session.trust_env = False
    current = url
    for _ in range(MAX_REDIRECTS + 1):
        safe, reason = _check_http_url_safety(current)
        if not safe:
            return False, "", reason

        try:
            resp = session.get(
                current,
                headers=headers,
                timeout=(10, 20),
                allow_redirects=False,
                stream=True,
            )
        except Exception:
            return False, "", "request_failed"

        # Follow redirects manually so we can safety-check the destination URL.
        if resp.status_code in {301, 302, 303, 307, 308}:
            loc = (resp.headers.get("Location") or "").strip()
            if not loc:
                return False, "", "request_failed"
            current = requests.compat.urljoin(current, loc)
            continue

        try:
            resp.raise_for_status()
        except Exception:
            return False, "", "request_failed"

        ctype = (resp.headers.get("Content-Type") or "").lower()
        if ("text/html" not in ctype) and ("application/xhtml" not in ctype) and ("text/plain" not in ctype):
            return False, "", "unsupported_content_type"

        # Stream and cap bytes to avoid huge responses causing memory issues.
        raw = bytearray()
        try:
            for chunk in resp.iter_content(chunk_size=65536):
                if not chunk:
                    continue
                raw.extend(chunk)
                if len(raw) > MAX_WEB_RESPONSE_BYTES:
                    return False, "", "response_too_large"
        except Exception:
            return False, "", "request_failed"

        try:
            encoding = resp.encoding or "utf-8"
            html = raw.decode(encoding, errors="ignore")
        except Exception:
            html = raw.decode("utf-8", errors="ignore")

        text = extract_readable_text(html)
        break
    else:
        return False, "", "request_failed"

    if len(text) < 250:
        return False, "", "too_little_text"

    return True, text[:WEB_TEXT_PER_URL_CAP], "ok"


def pseudo_pages_from_web_text(text: str) -> int:
    """Map extracted web text to a page count for max-question sizing."""
    if not text:
        return 0
    return max(1, len(text) // WEB_CHARS_PER_PAGE)


def docx_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    doc = Document(input_path)
    c = canvas.Canvas(output_path)

    y = 800
    for para in doc.paragraphs:
        text = para.text
        if y < 50:
            c.showPage()
            y = 800
        c.drawString(50, y, text[:100])
        y -= 15

    c.save()
    return output_path


def pptx_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    prs = Presentation(input_path)
    c = canvas.Canvas(output_path)

    for slide in prs.slides:
        y = 800
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if y < 50:
                    c.showPage()
                    y = 800
                c.drawString(50, y, text[:100])
                y -= 15
        c.showPage()

    c.save()
    return output_path


def image_to_pdf(input_path: str) -> str:
    output_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pdf")

    image = Image.open(input_path).convert("RGB")
    image.save(output_path, "PDF")

    return output_path


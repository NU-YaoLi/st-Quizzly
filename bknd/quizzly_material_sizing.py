"""
Format-agnostic “material size” for capping the Number of Questions slider.

Maps uploaded files and website text chunks to pseudo-pages (same scale as
``WEB_CHARS_PER_PAGE`` in ``quizzly_config``), so PDF / DOCX / PPTX / images and
URLs are comparable. PDFs combine extracted text volume with physical page
count so scans still scale while dense single pages still qualify.
"""

from __future__ import annotations

import os

import PyPDF2

from quizzly_config import WEB_CHARS_PER_PAGE

_IMAGE_EXTS = frozenset({".png", ".jpg", ".jpeg"})
_STRUCTURE_EXTS = frozenset({".docx", ".pptx"})


def char_volume_to_pseudo_pages(text: str) -> int:
    """
    Map plain text length to pseudo-pages. Matches
    ``pseudo_pages_from_web_text`` in ``quizzly_question_upldprcs``.
    """
    if not text:
        return 0
    return max(1, len(text) // WEB_CHARS_PER_PAGE)


def _pdf_extract_text_and_pages(path: str) -> tuple[str, int]:
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        n = len(reader.pages)
        parts: list[str] = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                parts.append("")
        return "".join(parts), n


def _pseudo_pages_for_pdf(path: str) -> int:
    text, phys = _pdf_extract_text_and_pages(path)
    text_units = char_volume_to_pseudo_pages(text)
    return max(text_units, phys)


def _pseudo_pages_for_docx(path: str) -> int:
    from docx import Document

    doc = Document(path)
    body = "\n".join(p.text for p in doc.paragraphs if p.text)
    u = char_volume_to_pseudo_pages(body)
    return max(1, u)


def _pseudo_pages_for_pptx(path: str) -> int:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    body = "\n".join(parts)
    u = char_volume_to_pseudo_pages(body)
    return max(1, u)


def estimate_upload_pseudo_pages(file_path: str) -> int:
    """
    One uploaded file (original path on disk, with correct extension).
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        try:
            return _pseudo_pages_for_pdf(file_path)
        except Exception:
            return 1

    if ext in _STRUCTURE_EXTS:
        try:
            if ext == ".docx":
                return _pseudo_pages_for_docx(file_path)
            return _pseudo_pages_for_pptx(file_path)
        except Exception:
            return 1

    if ext in _IMAGE_EXTS:
        return 1

    return 1


def total_pseudo_pages_for_upload_paths(paths: list[str]) -> int:
    return sum(estimate_upload_pseudo_pages(p) for p in paths)


def total_pseudo_pages_for_web_texts(text_chunks: list[str]) -> int:
    """Sum of pseudo-pages per fetched URL block (same rule as legacy UI)."""
    if not text_chunks:
        return 0
    return sum(char_volume_to_pseudo_pages(t) for t in text_chunks)
